from pptx import Presentation
from docx import Document
from io import BytesIO
from docx.shared import Inches, Pt
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

def extract_text_from_shape(shape):
    """Extract text from a shape."""
    if hasattr(shape, 'text'):
        return shape.text
    return ''

def extract_images_from_slide(slide, doc):
    """Extract images from a slide and add them to the document."""
    for shape in slide.shapes:
        if shape.shape_type == 13:  # Picture shape type
            image = shape.image
            image_bytes = image.blob

            img_stream = BytesIO(image_bytes)
            doc.add_picture(img_stream, width=Inches(4))
            doc.add_paragraph("Figure: Image from the slide").alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

def convert_pptx_to_docx(input_pptx, output_docx):
    """Convert a PowerPoint presentation to a Word document."""
    presentation = Presentation(input_pptx)
    doc = Document()

    for slide_number, slide in enumerate(presentation.slides):
        # Insert heading for each slide
        heading = f"Slide {slide_number + 1}"
        doc.add_heading(heading, level=1).alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

        # Extract text
        text = '\n'.join(extract_text_from_shape(shape) for shape in slide.shapes)
        doc.add_paragraph(text, style='BodyText')

        # Extract and insert images
        extract_images_from_slide(slide, doc)

    # Save the Word document
    doc.save(output_docx)

# Example usage
input_pptx = r''
output_docx = r''

convert_pptx_to_docx(input_pptx, output_docx)
